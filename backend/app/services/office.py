from __future__ import annotations

import os
import shutil
import tempfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pypdf import PdfReader

from .office_word import inspect_docx, mutate_docx
from .workspace import project_root, safe_path


EDITABLE_SUFFIXES = {".docx", ".xlsx", ".pptx"}


def inspect(project: dict, relative: str) -> dict:
    path = safe_path(project, relative)
    return _inspect_path(path)


def _inspect_path(path: Path) -> dict:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return inspect_docx(path)
    if suffix == ".xlsx":
        book = load_workbook(path, read_only=False, data_only=False)
        sheets = []
        for sheet in book.worksheets:
            rows = []
            formulas = []
            for row in list(sheet.iter_rows())[:200]:
                values = []
                for cell in row[:100]:
                    values.append(cell.value)
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas.append({"cell": cell.coordinate, "formula": cell.value})
                rows.append(values)
            sheets.append(
                {
                    "name": sheet.title,
                    "rows": rows,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                    "formulas": formulas[:500],
                    "merged_ranges": [str(value) for value in sheet.merged_cells.ranges],
                    "freeze_panes": str(sheet.freeze_panes or ""),
                    "auto_filter": sheet.auto_filter.ref or "",
                }
            )
        defined_names = [name for name in book.defined_names]
        book.close()
        return {"kind": "excel", "sheets": sheets, "defined_names": defined_names}
    if suffix == ".pptx":
        deck = Presentation(path)
        slides = []
        for index, slide in enumerate(deck.slides, 1):
            shapes = []
            for shape in slide.shapes:
                item = {"name": shape.name, "type": str(shape.shape_type)}
                if hasattr(shape, "text") and shape.text:
                    item["text"] = shape.text
                shapes.append(item)
            slides.append({"number": index, "shapes": shapes, "text": [item["text"] for item in shapes if item.get("text")]})
        return {"kind": "powerpoint", "slides": slides, "slide_count": len(slides)}
    if suffix == ".pdf":
        reader = PdfReader(path)
        return {
            "kind": "pdf",
            "pages": [{"number": i + 1, "text": (page.extract_text() or "")[:20_000]} for i, page in enumerate(reader.pages)],
            "page_count": len(reader.pages),
        }
    raise ValueError("Supported Office formats are DOCX, XLSX, PPTX and PDF")


def _add_word_content(doc: Document, content: str) -> None:
    for raw in content.splitlines():
        line = raw.rstrip()
        if not line:
            doc.add_paragraph("")
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)


def _format_sheet(sheet) -> None:
    if sheet.max_row:
        for cell in sheet[1]:
            cell.font = Font(bold=True)
        sheet.freeze_panes = "A2"
        if sheet.max_column:
            sheet.auto_filter.ref = f"A1:{get_column_letter(sheet.max_column)}{sheet.max_row}"
    for column in range(1, min(sheet.max_column, 50) + 1):
        letter = get_column_letter(column)
        width = 10
        for cell in list(sheet[letter])[:200]:
            width = max(width, min(len(str(cell.value or "")) + 2, 40))
        sheet.column_dimensions[letter].width = width


def create(project: dict, kind: str, relative: str, title: str, content: str, data: list[Any]) -> dict:
    if kind in {"preview", "edit"}:
        if not all(isinstance(item, dict) for item in data):
            raise ValueError("Office edit data must contain structured operation objects")
        operations = [dict(item) for item in data]
        return preview_edit(project, relative, operations) if kind == "preview" else edit(project, relative, operations)

    path = safe_path(project, relative)
    path.parent.mkdir(parents=True, exist_ok=True)
    if kind == "docx":
        doc = Document()
        doc.add_heading(title, 0)
        _add_word_content(doc, content)
        doc.save(path)
    elif kind == "xlsx":
        book = Workbook()
        sheet = book.active
        sheet.title = title[:31] or "Sheet1"
        for row in data or [[title], [content]]:
            if not isinstance(row, list):
                raise ValueError("Excel creation data must contain rows")
            sheet.append(row)
        _format_sheet(sheet)
        book.save(path)
    elif kind == "pptx":
        deck = Presentation()
        sections = [part.strip() for part in content.split("\n---\n") if part.strip()] or [content]
        for index, section in enumerate(sections):
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            lines = section.splitlines()
            slide.shapes.title.text = title if index == 0 else (lines[0][:120] if lines else f"{title} {index + 1}")
            body_lines = lines if index == 0 else lines[1:]
            slide.placeholders[1].text = "\n".join(body_lines)
        deck.save(path)
    else:
        raise ValueError("Unsupported Office output")
    return {"path": relative, "kind": kind, "size": path.stat().st_size}


def preview_edit(project: dict, relative: str, operations: list[dict[str, Any]]) -> dict:
    source = _editable_path(project, relative)
    before = _inspect_path(source)
    with tempfile.TemporaryDirectory(prefix="olladex-office-preview-") as directory:
        preview_path = Path(directory) / source.name
        shutil.copy2(source, preview_path)
        _mutate(project, preview_path, operations)
        after = _inspect_path(preview_path)
    return {
        "status": "preview",
        "path": _normalized_relative(project, source),
        "operation_count": len(operations),
        "operations": operations,
        "before": before,
        "after": after,
    }


def edit(project: dict, relative: str, operations: list[dict[str, Any]]) -> dict:
    source = _editable_path(project, relative)
    normalized = _normalized_relative(project, source)
    before = _inspect_path(source)
    temporary = source.with_name(f".{source.stem}.olladex-edit-{os.getpid()}{source.suffix}")
    shutil.copy2(source, temporary)
    try:
        _mutate(project, temporary, operations)
        after = _inspect_path(temporary)
        backup = _backup_file(project, source, normalized)
        os.replace(temporary, source)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise
    return {
        "status": "applied",
        "path": normalized,
        "operation_count": len(operations),
        "operations": operations,
        "backup_path": backup,
        "size": source.stat().st_size,
        "before": before,
        "after": after,
    }


def _editable_path(project: dict, relative: str) -> Path:
    path = safe_path(project, relative)
    if not path.is_file():
        raise ValueError("Office file not found")
    if path.suffix.lower() not in EDITABLE_SUFFIXES:
        raise ValueError("Structured editing supports DOCX, XLSX and PPTX")
    return path


def _normalized_relative(project: dict, path: Path) -> str:
    return path.relative_to(project_root(project)).as_posix()


def _backup_file(project: dict, source: Path, normalized: str) -> str:
    stamp = datetime.now(UTC).strftime("%Y%m%dT%H%M%S.%fZ")
    backup_root = project_root(project) / ".olladex" / "history" / stamp
    backup = backup_root / normalized
    backup.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, backup)
    return backup.relative_to(project_root(project)).as_posix()


def _mutate(project: dict, path: Path, operations: list[dict[str, Any]]) -> None:
    if not operations:
        raise ValueError("At least one Office edit operation is required")
    suffix = path.suffix.lower()
    if suffix == ".docx":
        mutate_docx(project, path, operations)
    elif suffix == ".xlsx":
        _mutate_xlsx(path, operations)
    elif suffix == ".pptx":
        _mutate_pptx(path, operations)
    else:
        raise ValueError("Structured editing supports DOCX, XLSX and PPTX")


def _mutate_xlsx(path: Path, operations: list[dict[str, Any]]) -> None:
    book = load_workbook(path, read_only=False, data_only=False)
    try:
        for operation in operations:
            action = operation.get("action")
            if action == "set_cell":
                sheet = _sheet(book, operation.get("sheet"))
                cell = str(operation.get("cell", "")).strip().upper()
                if not cell:
                    raise ValueError("set_cell requires a cell address")
                sheet[cell] = operation.get("value")
            elif action == "add_sheet":
                name = str(operation.get("name", "")).strip()
                if not name:
                    raise ValueError("add_sheet requires a name")
                if name in book.sheetnames:
                    raise ValueError(f"Worksheet already exists: {name}")
                book.create_sheet(title=name[:31])
            elif action == "rename_sheet":
                sheet = _sheet(book, operation.get("sheet"))
                name = str(operation.get("name", "")).strip()
                if not name:
                    raise ValueError("rename_sheet requires a name")
                sheet.title = name[:31]
            else:
                raise ValueError(f"Unsupported XLSX edit action: {action}")
        book.save(path)
    finally:
        book.close()


def _mutate_pptx(path: Path, operations: list[dict[str, Any]]) -> None:
    deck = Presentation(path)
    for operation in operations:
        action = operation.get("action")
        if action == "set_shape_text":
            slide_index = _index(operation, "slide_index", len(deck.slides))
            slide = deck.slides[slide_index]
            shape_index = _index(operation, "shape_index", len(slide.shapes))
            shape = slide.shapes[shape_index]
            if not hasattr(shape, "text"):
                raise ValueError("Selected PowerPoint shape does not contain editable text")
            shape.text = str(operation.get("text", ""))
        elif action == "add_slide":
            layout_index = int(operation.get("layout_index", 1))
            if layout_index < 0 or layout_index >= len(deck.slide_layouts):
                raise ValueError("PowerPoint layout_index is out of range")
            slide = deck.slides.add_slide(deck.slide_layouts[layout_index])
            title = str(operation.get("title", ""))
            content = str(operation.get("content", ""))
            if slide.shapes.title is not None:
                slide.shapes.title.text = title
            for placeholder in slide.placeholders:
                if placeholder == slide.shapes.title:
                    continue
                if hasattr(placeholder, "text"):
                    placeholder.text = content
                    break
        else:
            raise ValueError(f"Unsupported PPTX edit action: {action}")
    deck.save(path)


def _index(operation: dict[str, Any], key: str, length: int) -> int:
    try:
        index = int(operation.get(key))
    except (TypeError, ValueError) as exc:
        raise ValueError(f"{key} must be an integer") from exc
    if index < 0 or index >= length:
        raise ValueError(f"{key} is out of range")
    return index


def _sheet(book, requested: Any):
    name = str(requested or "").strip()
    if not name:
        return book.active
    if name not in book.sheetnames:
        raise ValueError(f"Worksheet not found: {name}")
    return book[name]
