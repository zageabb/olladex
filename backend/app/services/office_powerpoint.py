from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .workspace import safe_path


SHAPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "hexagon": MSO_SHAPE.HEXAGON,
}


def inspect_pptx(path: Path) -> dict:
    deck = Presentation(path)
    slides = []
    for slide_index, slide in enumerate(deck.slides):
        shapes = []
        for shape_index, shape in enumerate(slide.shapes):
            item = {
                "index": shape_index,
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_inches": round(shape.left.inches, 3),
                "top_inches": round(shape.top.inches, 3),
                "width_inches": round(shape.width.inches, 3),
                "height_inches": round(shape.height.inches, 3),
                "rotation": float(shape.rotation or 0),
                "text": shape.text if hasattr(shape, "text") else "",
                "placeholder": bool(getattr(shape, "is_placeholder", False)),
            }
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraphs.append(
                        {
                            "text": paragraph.text,
                            "level": paragraph.level,
                            "alignment": str(paragraph.alignment or ""),
                            "runs": [
                                {
                                    "text": run.text,
                                    "bold": bool(run.font.bold),
                                    "italic": bool(run.font.italic),
                                    "font_size": run.font.size.pt if run.font.size else None,
                                }
                                for run in paragraph.runs
                            ],
                        }
                    )
                item["paragraphs"] = paragraphs
            fill = _shape_fill(shape)
            line = _shape_line(shape)
            if fill:
                item["fill_color"] = fill
            if line:
                item["line_color"] = line
            shapes.append(item)
        slides.append(
            {
                "number": slide_index + 1,
                "index": slide_index,
                "layout": slide.slide_layout.name,
                "shapes": shapes,
                "text": [item["text"] for item in shapes if item.get("text")],
            }
        )
    return {
        "kind": "powerpoint",
        "slides": slides,
        "slide_count": len(slides),
        "slide_width_inches": round(deck.slide_width.inches, 3),
        "slide_height_inches": round(deck.slide_height.inches, 3),
        "layouts": [{"index": index, "name": layout.name} for index, layout in enumerate(deck.slide_layouts)],
    }


def mutate_pptx(project: dict, path: Path, operations: list[dict[str, Any]]) -> None:
    deck = Presentation(path)
    for operation in operations:
        action = operation.get("action")
        if action == "set_shape_text":
            shape = _shape(deck, operation)
            if not hasattr(shape, "text"):
                raise ValueError("Selected PowerPoint shape does not contain editable text")
            shape.text = str(operation.get("text", ""))
        elif action == "set_shape_position":
            shape = _shape(deck, operation)
            if "left" in operation:
                shape.left = Inches(float(operation["left"]))
            if "top" in operation:
                shape.top = Inches(float(operation["top"]))
            if "width" in operation:
                shape.width = Inches(float(operation["width"]))
            if "height" in operation:
                shape.height = Inches(float(operation["height"]))
            if "rotation" in operation:
                shape.rotation = float(operation["rotation"])
        elif action == "set_shape_style":
            shape = _shape(deck, operation)
            if "fill_color" in operation:
                color = _hex(operation.get("fill_color"))
                if color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(color)
                else:
                    shape.fill.background()
            if "line_color" in operation:
                color = _hex(operation.get("line_color"))
                if color:
                    shape.line.color.rgb = RGBColor.from_string(color)
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "font_size" in operation:
                            run.font.size = Pt(float(operation["font_size"]))
                        if "bold" in operation:
                            run.font.bold = bool(operation["bold"])
                        if "italic" in operation:
                            run.font.italic = bool(operation["italic"])
                        if "font_color" in operation:
                            color = _hex(operation.get("font_color"))
                            if color:
                                run.font.color.rgb = RGBColor.from_string(color)
        elif action == "add_slide":
            layout_index = int(operation.get("layout_index", 1))
            if layout_index < 0 or layout_index >= len(deck.slide_layouts):
                raise ValueError("PowerPoint layout_index is out of range")
            slide = deck.slides.add_slide(deck.slide_layouts[layout_index])
            title = str(operation.get("title", ""))
            content = str(operation.get("content", ""))
            if slide.shapes.title is not None:
                slide.shapes.title.text = title
            for placeholder in slide.placeholders:
                if placeholder == slide.shapes.title:
                    continue
                if hasattr(placeholder, "text"):
                    placeholder.text = content
                    break
        elif action == "delete_slide":
            slide_index = _index(operation, "slide_index", len(deck.slides))
            slide_id = deck.slides._sldIdLst[slide_index]
            relationship_id = slide_id.rId
            deck.part.drop_rel(relationship_id)
            deck.slides._sldIdLst.remove(slide_id)
        elif action == "reorder_slide":
            slide_index = _index(operation, "slide_index", len(deck.slides))
            target_index = int(operation.get("target_index", slide_index))
            if target_index < 0 or target_index >= len(deck.slides):
                raise ValueError("target_index is out of range")
            slide_id = deck.slides._sldIdLst[slide_index]
            deck.slides._sldIdLst.remove(slide_id)
            deck.slides._sldIdLst.insert(target_index, slide_id)
        elif action == "add_textbox":
            slide = _slide(deck, operation)
            textbox = slide.shapes.add_textbox(
                Inches(float(operation.get("left", 1))),
                Inches(float(operation.get("top", 1))),
                Inches(float(operation.get("width", 4))),
                Inches(float(operation.get("height", 1))),
            )
            textbox.text = str(operation.get("text", "Text"))
        elif action == "add_image":
            slide = _slide(deck, operation)
            image_path = str(operation.get("image_path", "")).strip()
            if not image_path:
                raise ValueError("add_image requires image_path")
            image = safe_path(project, image_path)
            if not image.is_file():
                raise ValueError("Image file not found in project")
            width = operation.get("width")
            height = operation.get("height")
            slide.shapes.add_picture(
                str(image),
                Inches(float(operation.get("left", 1))),
                Inches(float(operation.get("top", 1))),
                width=Inches(float(width)) if width else None,
                height=Inches(float(height)) if height else None,
            )
        elif action == "add_shape":
            slide = _slide(deck, operation)
            shape_name = str(operation.get("shape", "rectangle")).lower()
            if shape_name not in SHAPES:
                raise ValueError(f"Unsupported PowerPoint shape: {shape_name}")
            shape = slide.shapes.add_shape(
                SHAPES[shape_name],
                Inches(float(operation.get("left", 1))),
                Inches(float(operation.get("top", 1))),
                Inches(float(operation.get("width", 2))),
                Inches(float(operation.get("height", 1))),
            )
            if hasattr(shape, "text"):
                shape.text = str(operation.get("text", ""))
            color = _hex(operation.get("fill_color"))
            if color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(color)
        elif action == "set_slide_background":
            slide = _slide(deck, operation)
            color = _hex(operation.get("color"))
            if not color:
                raise ValueError("set_slide_background requires color")
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(color)
        else:
            raise ValueError(f"Unsupported PPTX edit action: {action}")
    deck.save(path)


def _slide(deck: Presentation, operation: dict[str, Any]):
    index = _index(operation, "slide_index", len(deck.slides))
    return deck.slides[index]


def _shape(deck: Presentation, operation: dict[str, Any]):
    slide = _slide(deck, operation)
    index = _index(operation, "shape_index", len(slide.shapes))
    return slide.shapes[index]


def _index(operation: dict[str, Any], key: str, length: int) -> int:
    try:
        index = int(operation.get(key))
    except (TypeError, ValueError) as exc:
        raise ValueError(f"{key} must be an integer") from exc
    if index < 0 or index >= length:
        raise ValueError(f"{key} is out of range")
    return index


def _hex(value: Any) -> str:
    text = str(value or "").strip().lstrip("#").upper()
    if not text:
        return ""
    if len(text) != 6 or any(character not in "0123456789ABCDEF" for character in text):
        raise ValueError("Colour must be a 6 digit hexadecimal value")
    return text


def _shape_fill(shape) -> str:
    try:
        if shape.fill.type is None:
            return ""
        rgb = shape.fill.fore_color.rgb
        return str(rgb) if rgb else ""
    except (AttributeError, TypeError):
        return ""


def _shape_line(shape) -> str:
    try:
        rgb = shape.line.color.rgb
        return str(rgb) if rgb else ""
    except (AttributeError, TypeError):
        return ""
